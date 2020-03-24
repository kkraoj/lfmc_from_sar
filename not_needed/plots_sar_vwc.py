# -*- coding: utf-8 -*-
"""
Created on Tue May 29 11:27:49 2018

@author: kkrao
"""
import os
import pandas as pd
import numpy as np
import matplotlib.pyplot as plt
import seaborn as sns
import matplotlib.ticker as mtick
from scipy import stats
from scipy.stats import gaussian_kde
from matplotlib import rcParams, colors, cm
from mpl_toolkits.basemap import Basemap
from matplotlib.patches import Polygon
from mpl_toolkits.axes_grid1 import make_axes_locatable
from sklearn.metrics import mean_squared_error
from sklearn import linear_model
import matplotlib as mpl
import matplotlib.dates as mdates
from matplotlib.ticker import MaxNLocator
from pptx import Presentation 
from pptx.util import Inches
from pathlib import Path
import glob
species_groups = {
'sagebrush':['Sagebrush, Silver','Sagebrush, Mountain Big',  'Sagebrush, Basin Big',
            'Sagebrush, California', 'Sagebrush, Black','Sagebrush, Wyoming Big', 
            'Sagebrush, Sand','Sagebrush, Bigelows', 'Sage, Black'],
'pine': ['Pine, Ponderosa','Pine, Lodgepole','Pine, Interior Ponderosa',
            'Pine, Loblolly'],
'chamise':['Chamise, Old Growth','Chamise','Chamise, New Growth'],
'manzanita': ['Manzanita, Whiteleaf','Manzanita, Eastwoods','Manzanita, Pinemat',
        'Manzanita, Greenleaf','Manzanita, Pointleaf'],
'oak': ['Oak, Texas Red','Oak, Live','Oak, Gambel', 'Oak, Sonoran Scrub',
        'Oak, Water','Tanoak','Oak, California Live','Oak, Gray', 'Oak, Emory'],
'juniper': [ 'Juniper, Redberry','Juniper, Rocky Mountain',  'Juniper, Utah',
           'Juniper, Ashe','Juniper, Creeping','Juniper, Oneseed', 'Juniper, Alligator',
           'Juniper, Western'],
'ceonothus': ['Ceanothus, Whitethorn','Ceanothus, Bigpod', 'Ceanothus, Redstem',
           'Ceanothus, Desert', 'Ceanothus, Buckbrush','Ceanothus, Deerbrush',       
             'Ceanothus, Hoaryleaf',   'Ceanothus, Snowbrush'],
'fir' :['Fir, California Red','Douglas-Fir', 'Fir, Subalpine',   'Fir, Grand',
          'Douglas-Fir, Coastal','Douglas-Fir, Rocky Mountain',  'Fir, White'],
'others': ['Mesquite, Honey', 'Bitterbrush, Desert',
            'Red Shank','Pinyon, Twoneedle', 'Cedar, Incense', 'Pinyon, Mexican',
           'Pinyon, Singleleaf',  'Bitterbrush, Antelope',
           'Buckwheat, Eastern Mojave', 'Snowberry, Mountain',
            'Spruce, Engelmann', 'Chinquapin, Bush',
            'Tamarisk', 'Sage, Purple',
           'Coyotebrush', 'Redcedar, Eastern',  'Forage kochia',
           'Snowberry, Western', 'Fescue, Arizona',  'Maple, Rocky Mountain',
           'Yaupon', 'Duff (DC)','Bluestem, Little',
           'Pinegrass',  'Sumac, Evergreen',  'Ninebark, Pacific']
}
def clean_xy(x,y,rep_times=1,thresh=0.0):
    # for testing ONLY
#    x=data_anomaly.values.flatten()
#    y=np.log10(mort.values.flatten())
    non_nan_ind=np.where(~np.isnan(x))[0]
    x=x.take(non_nan_ind);y=y.take(non_nan_ind)
    non_nan_ind=np.where(~np.isnan(y))[0]
    x=x.take(non_nan_ind);y=y.take(non_nan_ind)
    inds=np.where(x>=thresh)[0]
    x=x.take(inds)  ;y=y.take(inds)
    x=np.repeat(x,rep_times);y=np.repeat(y,rep_times)
    xy = np.vstack([x,y])
    z = gaussian_kde(xy)(xy)
    idx = z.argsort()
    x, y, z = x[idx], y[idx], z[idx]
#    x, y, z = np.reshape(x,(len(x),1)), np.reshape(y,(len(y),1)),\
#                        np.reshape(z,(len(z),1))
    return x,y,z

def adjust_spines(ax, spines):
    for loc, spine in ax.spines.items():
        if loc in spines:
            spine.set_position(('outward', 10))  # outward by 10 points
            spine.set_smart_bounds(True)
        else:
            spine.set_color('none')  # don't draw spine

    # turn off ticks where there is no spine
    if 'left' in spines:
        ax.yaxis.set_ticks_position('left')
#    else:
        # no yaxis ticks
#        ax.yaxis.set_ticks([])

    if 'bottom' in spines:
        ax.xaxis.set_ticks_position('bottom')

def normalize(df, param, norm):
    minmax = pd.read_pickle('data/df_sar')
#    minmax.groupby('Site').VH.min()
#    df.index = df.Site
#    temp = ((df[norm] - minmax.groupby(param)[norm].min())/\
#         (minmax.groupby(param)[norm].max()- minmax.groupby(param)[norm].min()))
#    temp.dropna()

    
    df=df.copy()
##    df.loc[:,norm] -= df.loc[:,'VV'].max()
    for p in df[param].unique():
##        df.loc[df[param]==p,norm]-=df.loc[df[param]==p,norm].max()
#        mn = minmax.loc[minmax[param]==p,norm].min()
#        mx = minmax.loc[minmax[param]==p,norm].max()
        mn = df.loc[df[param]==p,norm].min()
        mx = df.loc[df[param]==p,norm].max()
        
        df.loc[df[param]==p,norm] = (df.loc[df[param]==p,norm]-mn)/\
        (mx - mn)
    return df




def transform(df, totransform='VH', basedon="Fuel"):
    def minmax(df):
        mn = df[totransform].min()
        mx = df[totransform].max()
        new_col_name = "%s_%s"%(totransform, basedon)
        df[new_col_name] = (df[totransform]-mn)/(mx-mn)
        return df
    df = df.groupby(basedon).apply(minmax)  
    return df

def plot_usa(enlarge = 1.): 
    sns.set_style('ticks')
    fig, ax = plt.subplots(figsize=(8*enlarge,5*enlarge))
    m = Basemap(llcrnrlon=-119,llcrnrlat=22,urcrnrlon=-64,urcrnrlat=49,
            projection='lcc',lat_1=33,lat_2=45,lon_0=-95)
    m = Basemap(llcrnrlon=-123,llcrnrlat=23,urcrnrlon=-90,urcrnrlat=50,
            projection='lcc',lat_1=33,lat_2=45,lon_0=-105)
    m.drawmapboundary(fill_color='lightcyan')
    #-----------------------------------------------------------------------
    # load the shapefile, use the name 'states'
    m.readshapefile('D:/Krishna/projects/vwc_from_radar/cb_2017_us_state_500k', 
                    name='states', drawbounds=True)
    statenames=[]
    for shapedict in m.states_info:
        statename = shapedict['NAME']
        statenames.append(statename)
    for nshape,seg in enumerate(m.states):
        if statenames[nshape] == 'Alaska':
        # Alaska is too big. Scale it down to 35% first, then transate it. 
            new_seg = [(0.35*args[0] + 1100000, 0.35*args[1]-1500000) for args in seg]
            seg = new_seg    
        poly = Polygon(seg,facecolor='papayawhip',edgecolor='k', zorder  = 1)
        ax.add_patch(poly)

    
    #, ticks = np.linspace(0,20,5)
    return fig, ax, m
def map_list_to_colors(landcover):
    cmap = plt.get_cmap('Accent')
    colors = cmap(np.linspace(0, 1, len(landcover.unique())))
    colors = dict(zip(landcover.unique(), colors))
    colors = landcover.map(colors)
    return colors
pd.set_option('display.max_columns', 20)
pd.set_option('display.max_rows', 100)
pd.set_option('display.width', 200)
sns.set(font_scale = 1.2, style = 'ticks')
rcParams['font.family'] = 'serif'
#rcParams['font.sans-serif'] = ['Tahoma']
def plot_heatmap(corr, method = 'spearman', x_name =  "Vs. Fuel Moisture" ,\
                 y_size = 3):
    clustergrid=sns.clustermap(corr,figsize=(1e-1,1e-1))
    rows = clustergrid.dendrogram_row.reordered_ind 
    cols = clustergrid.dendrogram_col.reordered_ind 
    corr=corr.iloc[rows,cols]
    corr.columns = corr.columns.str.upper()  
    corr.columns.name  =  x_name                
    fig, ax = plt.subplots(figsize = (5,5),dpi = 300)
    plt.yticks(fontsize=y_size)
    sns.heatmap(corr, ax = ax, vmin = 0, vmax = 1)


def make_patch_spines_invisible(ax):
    ax.set_frame_on(True)
    ax.patch.set_visible(False)
    for sp in ax.spines.values():
        sp.set_visible(False)
        
def export_images_to_ppt(imagepath, pptfile, imageformat = "png", \
                         ht = 2.5, wd = 4.3, lt = 1, tp = 0.3):
    images = glob.glob(imagepath+"\\*.%s"%imageformat)
    prs = Presentation(pptfile)
    blank_slide_layout = prs.slide_layouts[8] 
    for i, image in enumerate(images):
        if i%4==0:
            slide = prs.slides.add_slide(blank_slide_layout) 
            left = Inches(lt); top = Inches(tp)
            slide.shapes.add_picture(image, left, top, height = Inches(ht)) 
        elif i%4==1:
            left = Inches(lt)+Inches(wd); top = Inches(tp)
            slide.shapes.add_picture(image, left, top, height = Inches(ht)) 
        elif i%4==2:
            left = Inches(lt); top = Inches(tp)+Inches(ht)
            slide.shapes.add_picture(image, left, top, height = Inches(ht))     
        elif i%4==3:
            left = Inches(lt)+Inches(wd); top = Inches(tp)+Inches(ht)
            slide.shapes.add_picture(image, left, top, height = Inches(ht))                 
    prs.save(pptfile)
    
        
os.chdir('D:/Krishna/projects/vwc_from_radar')
df = pd.read_pickle('data/df_all')
#df.obs_date_local = pd.to_datetime(df.obs_date_local, format = '%Y-%m-%d %H:%M')
###############################################################################
####### histogram of residuals
#df = df.loc[(df.residual.abs() <=2),:]
#print(df.shape)
#fig, ax = plt.subplots()
#
#df.hist(column = 'residual', bins = 800, facecolor='g',edgecolor='black',  alpha=0.75, ax = ax)
#ax.set_xlabel('$\Delta$ delay (days)')
#ax.set_ylabel('Frequency (-)')
#ax.set_title('')
#ax.grid('off')
##ax.set_xlim([-14,14])
##ax.axvline(x=-5, ls = '--', c='r')
##ax.axvline(x=5, ls = '--', c='r')

###############################################################################
####### plot of vwc vs sar regression scatter
#df=normalize(df, 'Fuel', 'VH')
#filter = (df.residual.abs() <=6)
#d = df.loc[filter,:]
##filter = (df.residual.abs() <=2)&(df.Site == 'Deer Hill')
##d = df.loc[filter,:]
#fig, ax = plt.subplots(figsize = (3,3))
#cmap = "magma"
#response = "VH"
#predictor = "Percent"
#x,y,z=clean_xy(d[predictor].values, d[response].values, thresh = -100, rep_times = 1)
#ax.scatter(x,y,c=z,edgecolor='',cmap=cmap)
#ax.xaxis.set_major_formatter(mtick.PercentFormatter())
#ax.set_ylabel(r'$\frac{\sigma_{VH} - min_{species}\sigma_{VH}}{max_{species}\sigma_{VH} - min_{species}\sigma_{VH}}$')
##ax.set_xlabel('$\sigma_{%s}$ (dB)'%predictor)
#ax.set_xlabel(r'Fuel Moisture')
##ax.set_ylim([-20,0])
##ax.set_ylim([-30,-5])
#adjust_spines(ax,['left', 'bottom'])
#ax.set_xlim([0,300])
#r2 = d[predictor].corr(d[response])**2
#ax.annotate('$R^2$ = %0.2f'%r2, xy=(0.65, 0.89), xycoords='axes fraction')
###############################################################################

##### summary table of R squareds for sites
##### something is not working  for normalizing VH
#df=normalize(df, 'Site', 'VH')
#filter = (df.residual.abs() <=6)&(~df.Fuel.isin(['10-Hour', '100-Hour', '1000-Hour']))
#groupby =["Site", 'Fuel']
##print(df.VH.min(), df.VH.max())
#d = df.loc[filter,:].groupby(groupby).filter(lambda x: len(x)>=20)
#
#d = pd.DataFrame(d.groupby(groupby).apply(lambda x: \
#                 np.array(stats.linregress(x['Percent'], x['VH']))\
#                           [[2,3]]), columns = ['p_value'])
#d['r2'] = [x[0]**2 for x in d['p_value']]
#d['p_value'] = [x[1] for x in d['p_value']]
#d.dropna(inplace = True)
##filters = (d.r2 >= 0.1)&(d.p_value<=0.05)
##d = d.loc[filters,:]
#print(d.sort_values("p_value"))
###############################################################################
########## time series of vwc and sar
#predictor = "VH_corr" 
#i=0
#for site in df.Site.unique():
#    filter =\
#            (df.Site == site)\
#            &(df.residual.abs() <=6)\
#            &(~df.Fuel.isin(['10-Hour', '100-Hour', '1000-Hour']))\
#            &(df.ndvi>=0.2)
##    for fuel in df.loc[filter,'Fuel'].unique():
##    d = df.loc[filter&(df.Fuel == fuel),:]
#    d = df.loc[filter,:]
#    if len(d.obs_date_local.unique()) <10:
#        continue
#    g = d.groupby('meas_date')
#    fig, ax = plt.subplots(figsize = (7,3) )
#    ax.plot(g.obs_date_local.min(), g.Percent.mean(), 'b.-')
#    ax.set_ylabel('Fuel Moisture (%)', color='b')
#    ax.tick_params('y', colors='b')
#    ax2 = ax.twinx()
#    ax2.plot(g.obs_date_local.min(), g[predictor].mean(), 'm.-', linewidth = 1)
#    ax2.plot(g.obs_date_local.min(), g["VH"].mean(), 'm.--', linewidth = 1)
##    ax2.set_ylim([-30,-5])
#    ax2.set_ylabel('$\sigma_{VH}$ (dB)', color = 'm')
#    ax2.tick_params('y', colors='m')
#    ax3 = ax.twinx()
#    ax3.spines["right"].set_position(("axes", 1.6))
#    make_patch_spines_invisible(ax3)
#    ax3.spines["right"].set_visible(True)
#    ax3.plot(g.obs_date_local.min(), g[predictor].mean()/g.ndvi.mean(), 'g.-', linewidth = 1)
#    ax3.set_ylabel(r'$\frac{\sigma_{VH} (dB)}{NDVI}$', color = 'g')
#    ax3.tick_params('y', colors='g')
#
#    fig.autofmt_xdate()
#    ax.xaxis.set_major_locator(MaxNLocator( 6))
#    ax.set_title(site)
#    os.chdir(Path(r"D:\Krishna\projects\vwc_from_radar\codes\plots"))
#    plt.tight_layout()
#    fig.savefig('%s.png'%i, bbox_inches = "tight")
#    i+=1
#    plt.show()
#export_images_to_ppt(imagepath = r"D:\Krishna\projects\vwc_from_radar\codes\plots",\
#    pptfile = 'C:\\Users\\kkrao\\Dropbox\\meetingsKrishna\\Presentations\\31-Jul-2018 (angle correction).pptx')
    
###############################################################################
## plot time series of SAR
#df = pd.read_pickle('data/df_sar_pm')
#for site in df.Site.unique():
#    filter =\
#            (df.Site == site)
#    d = df.loc[filter,:]
#    if d.shape[0] <100:
#        continue
#    fig, ax = plt.subplots(figsize = (3,3) )
#    ax.plot(d.obs_date, d['VH'], 'b.-', linewidth = 1)
#    ax.set_ylabel('Fuel Moisture (%)', color='b')
#    ax.tick_params('y', colors='b')
#
##    ax.set_ylim([-20,0])
#    ax.set_ylabel('$\sigma_{VH}$ (dB)', color = 'b')
#    fig.autofmt_xdate()
#
#    plt.show()
###############################################################################
###combining optical data
#df2 = pd.read_pickle('data/df_optical_vwc') 
#df = df.join(df2, rsuffix = "_opt")
### check if propoerly joined
#check_columns = ['GACC', 'State', 'Group', 'Site','Fuel', 'Percent', 'meas_date']
##for col in check_columns:
##    print((df[col]!=df['%s_opt'%col]).sum())
#df.drop(['%s_opt'%x for x in check_columns],axis = 1, inplace = True)
#df.rename(columns = {'B8':'nir'}, inplace = True)
#df.to_pickle('data/df_sar_opt_vwc')
#df['VH-VV'] = df['VH'] - df['VV'] 
####### plot of vwc vs optical scatter
#filter = (df.residual.abs() <=6)&(df.residual_opt.abs() <=6)
#d = df.loc[filter,:]
##filter = (df.residual.abs() <=2)&(df.Site == 'Deer Hill')
##d = df.loc[filter,:]
#fig, ax = plt.subplots(figsize = (2,2))
#cmap = "viridis"
#response = "ndwi"
#predictor = "VH"
#x,y,z=clean_xy(d[predictor].values, d[response].values, thresh = -100, rep_times = 1)
#ax.scatter(x,y,c=z,edgecolor='',cmap=cmap)
##ax.xaxis.set_major_formatter(mtick.PercentFormatter())
#ax.set_ylabel('%s'%response.upper())
##ax.set_ylabel(r'$\sigma_{VH} - \sigma_{VV}$ (dB)')
#ax.set_xlabel(r'$\sigma_{VH}$ (dB)')
##ax.set_ylim([-30,-5])
#adjust_spines(ax,['left', 'bottom'])
##ax.set_xlim([0,200])
#ax.set_xlim([-30,-5])
#r2 = d[predictor].corr(d[response])**2
#ax.annotate('$R^2$ = %0.2f'%r2, xy=(0.65, 0.89), xycoords='axes fraction')
###############################################################################
##### heathmap of R squard with all predictors groupby site and fuel
#    ax.annotate('%s $R^2$ '%method, xy=(1.2, 1.02), \
#                xycoords='axes fraction', ha = 'right')
#df = pd.read_pickle('data/df_sar_opt_vwc')
#df = transform(df)
#df = transform(df, basedon="Site")
#filter = (df.residual.abs() <=6)&\
#            (~df.Fuel.isin(['10-Hour', '100-Hour', '1000-Hour']))&\
#            (df.residual_opt.abs() <=6)
#df = df.loc[filter,:]            
#groupby =["Site", 'Fuel']
#keep_cols = ['Site','Fuel', 'Percent', 'VV', 'VH','VH_Site', \
#             'VH_Fuel', 'nir', 'ndvi', 'ndwi', 'nirv']
#d = df.loc[:,keep_cols]
#d = d.groupby(groupby).filter(lambda x: len(x)>=10)
#method = 'pearson'
#corr_with = "Percent"
##corr = d.groupby(groupby).apply(lambda x: x.corr(method)[corr_with][1:]**2)
#corr = d.groupby(groupby).apply(lambda x: x.corr(method)[corr_with].\
#                 drop(corr_with)**2)
#corr = d.groupby(groupby).apply(lambda x: x.corr(method)[corr_with].\
#                 drop(corr_with)**2)
#plot_heatmap(corr, method)
#data_points = pd.DataFrame(d.groupby(groupby).apply(lambda x: x.shape[0]),\
#                           columns = ["Data Points"])
#fig, ax = plt.subplots(figsize = (0.5,5),dpi = 300)
#sns.heatmap(data_points, ax = ax, cmap = "viridis")
#plt.tick_params(axis='y',which='both',left=False,labelleft = False)
#ax.set_ylabel("")

###### heathmap of R squard with all predictors groupby site
#def mean_percent(df):
#    df.Percent = df.Percent.mean()
#    return df
#df = pd.read_pickle('data/df_sar_opt_vwc')
#df = transform(df)
#df = transform(df, basedon="Site")
##df=normalize(df, 'Site', 'VH')
#filter = (df.residual.abs() <=6)&\
#            (~df.Fuel.isin(['1-Hour', '10-Hour', '100-Hour', '1000-Hour']))&\
#            (df.residual_opt.abs() <=6)
#df = df.loc[filter,:]            
#groupby =["Site", 'meas_date']
#df = df.groupby(groupby).apply(mean_percent)
#df.drop_duplicates(['Site','meas_date','Percent'], inplace = True)
#keep_cols = ['Site', 'Percent', 'VV', 'VH', 'VH_Site', \
#             'VH_Fuel', 'nir', 'ndvi', 'ndwi', 'nirv']
#d = df.loc[:,keep_cols]
#d = d.groupby("Site").filter(lambda x: len(x)>=10)
#corr = d.groupby("Site").apply(lambda x: x.corr(method)[corr_with].\
#                 drop(corr_with)**2)

#plot_heatmap(corr, method, x_name="Vs. Average Fuel Moisture", y_size=6)
####### plot usa map
#variable = "VH"
#corr = corr[variable]
#latlon = pd.read_csv('data/fuel_moisture/nfmd_spatial.csv', index_col = 0)
#latlon["corr_%s"%variable] = corr
#fig, ax, cax, m = plot_usa()
#cmap = 'magma'
#plot=m.scatter(latlon.longitude.values, latlon.latitude.values, 
#               s=40,c=latlon["corr_%s"%variable],cmap =cmap ,edgecolor = 'k',\
#                    marker='o',latlon = True, zorder = 2,\
#                    vmin = 0, vmax = 1)
#plt.setp(ax.spines.values(), color='w')
#divider = make_axes_locatable(ax)
#cax = divider.append_axes("right", size="5%", pad=0.08)
#fig.colorbar(plot,ax=ax,cax=cax)
#ax.set_title('$R^2(\sigma_{VH}, Fuel\ Moisture)$')
#plt.show()
###############################################################################
#### heatmap vh R squared and means of latent variable
    
#df = pd.read_pickle('data/df_all')
#df = df.infer_objects()
##df = transform(df)
##df = transform(df, basedon="Site")
#filter = (df.residual.abs() <=6)&\
#            (~df.Fuel.isin(['10-Hour', '100-Hour', '1000-Hour']))&\
#            (df.residual_opt.abs() <=6)&\
#            (df.latitude!=0)
#df = df.loc[filter,:]            
#groupby =["Site", 'Fuel']
#predictor = "VH_corr"
#keep_cols = ['Site','Fuel', 'Percent', predictor, \
#              'nir', 'ndvi', 'ndwi', 'nirv','slope', 'elevation', \
#              'canopy_height', 'forest_cover', 'silt', 'sand', 'clay']
#d = df.loc[:,keep_cols]
#d = d.groupby(groupby).filter(lambda x: len(x)>=10)
#method = 'pearson'
#corr_with = "Percent"
#corr = d.groupby(groupby).apply(lambda x: x.corr(method)[corr_with].\
#                 drop(corr_with)**2)
##plot_heatmap(corr, method)
#cols = [e for e in list(d.columns) if e not in ['Site','Fuel']]
##d[cols] = (d[cols]  - d[cols].min())/(d[cols].max() - d[cols].min())
#mean = d.groupby(groupby).apply(lambda x: x.mean())
#mean[predictor] = corr[predictor]
#mean.rename(columns={predictor: r"$R^2(\sigma_{%s}, FM)$"%predictor, "forest_cover": "landcover"}, inplace = True)
#landcover = mean['landcover']
#mean.drop(["Percent", 'landcover'], axis = 1, inplace = True)

#colors = map_list_to_colors(landcover)
#mpl.rcParams['ytick.labelsize'] = 8
#mean[r"1-$R^2(\sigma_{%s}, FM)$"%predictor] = 1- mean[r"$R^2(\sigma_{%s}, FM)$"%predictor]
#cg = sns.clustermap(mean,figsize=(5,5), row_cluster = False, method = "complete", metric = 'correlation',\
#               standard_scale =1, row_colors =colors, cmap="mako")
######## plot look angle as function of site

#fig, ax = plt.subplots()
#df = df.loc[df.obs_date_local.dt.hour>16,:]
##df.obs_date_local.dt.hour.hist()
#df.groupby("Site").angle.std().hist(align = "mid", ax = ax)
#ax.set_xlabel(r"S.D. of look angle per site($^o$)")
#ax.set_ylabel("Frequency")
###############################################################################
### how many points exist such that two species are recorded at some location 
#delta_days = 14
#filter =\
#    (df.residual.abs() <=delta_days)\
#    &(~df.Fuel.isin(['10-Hour', '100-Hour', '1000-Hour']))\
#    &(df.ndvi>=0.1)\
#    &(df.Percent<=1000)\
#    &(df.Percent>=40)\
#    &(df.residual_opt.abs()<=delta_days)
#df = df.loc[filter,:]
#### on same day?
#groupby = ['Site','meas_date']
#counter = 0
#def count_reps(x):
#    if len(x.Fuel.unique)>1:
#        x.Fuel.unique()
#def select_similar_fm(x):
#    thresh = 10
##    if 35550 is in x.index:
##        print(x)
#    if x.Percent.max() - x.Percent.min() <= thresh:
##        print(x.index)
#        return x.index
    
##reps = df.groupby(groupby).apply(lambda x: x.Percent if len(x.Fuel.unique())>1 else 0)
##range = df.groupby(groupby).apply(lambda x: x.Percent.max() - x.Percent.min())
#similar = df.groupby(groupby).apply(select_similar_fm)
##print(similar)
##reps.drop()
##print(reps.sum()/df.shape[0])
##range.unstack(level=0).plot(kind = 'line',legend = False,style = '.')
#
#df.loc[df.Site=='Ziegler',['Site','meas_date','Percent']]
#df.loc[35550,:]
#
#pct_thresh = 20
#to_select = pd.DataFrame()
#for site in df.Site.unique():
#    d = df.loc[df.Site==site,:]
#    for date in d.meas_date.unique():
#        x = d.loc[d.meas_date == date,:]
##        print(x)
#        if x.Percent.max() - x.Percent.min()<= pct_thresh:
#            to_select = to_select.append(x)
#to_select.head()
#to_select.shape
###############################################################################
##### plot usa map with locations colored by seasonality
#latlon = pd.read_csv('data/fuel_moisture/nfmd_spatial.csv', index_col = 0)
#seasonality_of_sites = pd.read_excel("data/fuel_moisture/seasonality_of_locations.xlsx", index_col = 0, usecols = 1)
#dict_seasonality = {1:'Early peak',2:'Late peak',0:'Not seasonal'}
#fig, ax, m = plot_usa()
#cmap = plt.get_cmap('RdYlGn')
#cNorm  = colors.Normalize(vmin=0, vmax=2)
#scalarMap = cm.ScalarMappable(norm=cNorm, cmap=cmap)
#for i in seasonality_of_sites.seasonality.unique():
#    sites = seasonality_of_sites.loc[seasonality_of_sites.seasonality==i].index
#    plot=m.scatter(latlon.loc[sites,'longitude'].values, latlon.loc[sites,'latitude'].values, 
#                   s=20,color=scalarMap.to_rgba(i) ,edgecolor = 'k',\
#                        marker='o',latlon = True, label = dict_seasonality[i], zorder = 2)
#plt.setp(ax.spines.values(), color='w')
##divider = make_axes_locatable(ax)
##cax = divider.append_axes("right", size="5%", pad=0.08)
##fig.colorbar(plot,ax=ax,cax=cax)
#ax.legend(facecolor= 'white', frameon = True, handletextpad = -0.5)
#ax.set_title('Seasonality of sites')
##### plot usa map with locations colored by speak doy
#latlon = pd.read_csv('data/fuel_moisture/nfmd_spatial.csv', index_col = 0)
#seasonality_of_sites = pd.read_excel("data/fuel_moisture/seasonality_of_locations.xlsx", index_col = 0)
#fig, ax, m = plot_usa()
#cmap = 'viridis'
#latlon = seasonality_of_sites.join(latlon)
#plot=m.scatter(latlon.longitude.values, latlon.latitude.values, 
#               s=25,c = latlon.peak_doy,edgecolor = 'white',vmin = 50, vmax = 250,\
#                marker='o',latlon = True, zorder = 2, cmap = cmap)
#m.scatter(latlon.loc[latlon.seasonality==0,'longitude'].values, \
#          latlon.loc[latlon.seasonality==0,'latitude'].values, 
#               s=25,color = 'grey',edgecolor = 'white',\
#                marker='o',latlon = True, zorder = 3)## grey out non seasonal points
#plt.setp(ax.spines.values(), color='w')
#divider = make_axes_locatable(ax)
#cax = divider.append_axes("right", size="5%", pad=0.08)
#fig.colorbar(plot,ax=ax,cax=cax)
#ax.set_title('Peak day of year')
################################################################################
#### which species don't have seasonality
#meas = pd.read_pickle("data/df_vwc_historic")
#filter =\
#    (~meas.fuel.isin(['10-Hour', '100-Hour', '1000-Hour']))\
#    &(meas.percent<=200)\
#    &(meas.percent>=40)
#meas = meas.loc[filter,:]
#seasonality_of_sites = pd.read_excel("data/fuel_moisture/seasonality_of_locations.xlsx", index_col = 0)
#for site in seasonality_of_sites.index:
##    print(meas.loc[meas.site==site,'fuel'].unique()[0])
#    fuel = meas.loc[meas.site==site,'fuel'].unique()[0]
#    group = [group if fuel in species  else 'others' for group, species in species_groups.items()][0]
#    seasonality_of_sites.loc[site,'species_group']=group
#seasonality_of_sites.loc[seasonality_of_sites.seasonality>1,'seasonality']=1
#seasonality_of_sites.groupby('species_group').seasonality.value_counts()




