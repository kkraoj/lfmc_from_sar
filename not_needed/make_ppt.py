# -*- coding: utf-8 -*-
"""
Created on Tue Jul 31 04:29:50 2018

@author: kkrao
"""


from pptx import Presentation 
from pptx.util import Inches
import os
from pathlib import Path
import glob
from dirs import dir_codes

os.chdir(Path(r"D:\Krishna\projects\vwc_from_radar\codes\plots"))

#
#img_path = 'trial.png'

#prs = Presentation(os.path.join(ppt_dir,ppt_name ))
#blank_slide_layout = prs.slide_layouts[6] 
#slide = prs.slides.add_slide(blank_slide_layout)
#
#left = top = Inches(1.5)
#pic = slide.shapes.add_picture(img_path, left, top) 
#prs.save(ppt_name)


def export_images_to_ppt(imagepath, pptfile, imageformat = "jpg", \
                         ht = 2.5, wd = 6.4, lt = 1, tp = 1):
    images = glob.glob(imagepath+"\\*.%s"%imageformat)
    prs = Presentation(pptfile)
    blank_slide_layout = prs.slide_layouts[8] 
    for i, image in enumerate(images):
        if i%4==0:
            slide = prs.slides.add_slide(blank_slide_layout) 
            left = Inches(lt); top = Inches(tp)
            slide.shapes.add_picture(image, left, top, height = Inches(ht)) 
        elif i%4==1:
            left = Inches(lt)+Inches(wd); top = Inches(tp)
            slide.shapes.add_picture(image, left, top, height = Inches(ht)) 
        elif i%4==2:
            left = Inches(lt); top = Inches(tp)+Inches(ht)
            slide.shapes.add_picture(image, left, top, height = Inches(ht))     
        elif i%4==3:
            left = Inches(lt)+Inches(wd); top = Inches(tp)+Inches(ht)
            slide.shapes.add_picture(image, left, top, height = Inches(ht))                 
    prs.save(pptfile)


imagepath = dir_codes +os.sep + "plots"
ppt_dir = r"C:\Users\kkrao\Dropbox\meetingsKrishna\Presentations"
ppt_name = "14-May-2019 (site predictions).pptx"
pptfile = ppt_dir + os.sep + ppt_name
export_images_to_ppt(imagepath, pptfile)