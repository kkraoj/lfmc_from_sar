# -*- coding: utf-8 -*-
"""
Created on Tue Jul 31 04:29:50 2018

@author: kkrao
"""


from pptx import Presentation 
from pptx.util import Inches
import os
from pathlib import Path

os.chdir(Path(r"D:\Krishna\projects\vwc_from_radar\codes\plots"))


img_path = 'trial.png'
ppt_dir = r"C:\Users\kkrao\Dropbox\meetingsKrishna\Presentations"
ppt_name = "14-May-2019 (site predictions).pptx"
prs = Presentation(os.path.join(ppt_dir,ppt_name ))
blank_slide_layout = prs.slide_layouts[6] 
slide = prs.slides.add_slide(blank_slide_layout)

left = top = Inches(1.5)
pic = slide.shapes.add_picture(img_path, left, top) 
prs.save(ppt_name)
